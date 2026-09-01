#!/usr/bin/env python3
"""
Demo deck generator for the AWS Agentic AI Hackathon (Future Vision) at LEAP.

Builds an attractive, colorful PPTX from a small JSON content file, so teams get a
consistent, presentable deck and only fill in their own content. It follows the
required deliverable structure: idea, team, problem with a measurable claim,
solution, architecture, live demo, impact, and what you would build next.

Usage:
    pip install python-pptx
    python build_deck.py deck-content.json demo.pptx

Edit deck-content.json (see deck-content.example.json) with your idea and team,
then run the command above. Open the generated demo.pptx and record your 3-minute demo.
"""

import json
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# LEAP-inspired palette: deep indigo, violet, with warm and teal accents.
INK = RGBColor(0x14, 0x10, 0x3A)      # near-black indigo (dark backgrounds)
INDIGO = RGBColor(0x3B, 0x2A, 0x8C)   # primary
VIOLET = RGBColor(0x6D, 0x28, 0xD9)   # secondary
ACCENT = RGBColor(0xF9, 0x73, 0x16)   # warm accent (orange)
TEAL = RGBColor(0x14, 0xB8, 0xA6)     # cool accent
PAPER = RGBColor(0xFB, 0xFA, 0xFF)    # light background
SLATE = RGBColor(0x33, 0x33, 0x40)    # body text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

W, H = Inches(13.333), Inches(7.5)  # 16:9


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6):
    """runs: list of (text, size, color, bold)."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (txt, size, color, bold) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = "Arial"
    return box


def _bullets(slide, x, y, w, h, items, size=20, color=SLATE):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        r = p.add_run()
        r.text = "•  " + item
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Arial"


def _content_header(slide, title, kicker="FUTURE VISION"):
    _bg(slide, PAPER)
    _rect(slide, 0, 0, W, Inches(1.25), INDIGO)
    _rect(slide, 0, Inches(1.25), W, Inches(0.08), ACCENT)
    _text(slide, Inches(0.6), Inches(0.12), Inches(12), Inches(0.4),
          [(kicker, 12, TEAL, True)])
    _text(slide, Inches(0.6), Inches(0.42), Inches(12), Inches(0.8),
          [(title, 30, WHITE, True)])


def title_slide(prs, c):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, INK)
    _rect(s, 0, Inches(3.05), W, Inches(0.10), ACCENT)
    _text(s, Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.5),
          [("AWS AGENTIC AI HACKATHON  ·  LEAP", 14, TEAL, True)])
    _text(s, Inches(0.8), Inches(3.2), Inches(11.7), Inches(2.2),
          [(c.get("idea", "Your Idea"), 54, WHITE, True),
           (c.get("tagline", ""), 22, RGBColor(0xC7, 0xC2, 0xF0), False)])
    _text(s, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.8),
          [("Team " + c.get("team_name", ""), 20, ACCENT, True)])


def team_slide(prs, c):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, "The Team")
    members = c.get("members", [])
    cols = 2 if len(members) > 3 else 1
    col_w = Inches(6.0)
    for i, m in enumerate(members):
        col = i % cols
        row = i // cols
        x = Inches(0.7) + col * (col_w + Inches(0.4))
        y = Inches(1.7) + row * Inches(1.0)
        _rect(s, x, y, Inches(0.12), Inches(0.7), VIOLET)
        org = m.get("org") or m.get("flag") or "Independent"
        _text(s, x + Inches(0.3), y - Inches(0.05), col_w, Inches(0.9),
              [(m.get("name", ""), 22, INK, True),
               (org, 15, SLATE, False)], space=2)


def problem_slide(prs, c):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, "The Problem")
    _text(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(1.6),
          [("Who it affects", 14, VIOLET, True),
           (c.get("beneficiary", ""), 24, INK, False)], space=4)
    _rect(s, Inches(0.7), Inches(3.9), Inches(11.9), Inches(1.8), RGBColor(0xF1, 0xEE, 0xFC))
    _text(s, Inches(1.0), Inches(4.1), Inches(11.3), Inches(1.4),
          [("The measurable claim", 14, ACCENT, True),
           (c.get("claim", ""), 26, INDIGO, True)], space=4)


def bullets_slide(prs, title, items):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, title)
    _bullets(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(5.2), items or [])
    return s


def architecture_slide(prs, c):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, "Architecture")
    note = c.get("architecture_note", "")
    img = c.get("arch_image")
    if img:
        try:
            s.shapes.add_picture(img, Inches(0.8), Inches(1.7), height=Inches(4.9))
            if note:
                _text(s, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.6),
                      [(note, 14, SLATE, False)])
            return s
        except Exception as exc:  # noqa: BLE001
            note = f"{note}  (add your diagram image: {exc})"
    _rect(s, Inches(0.8), Inches(1.9), Inches(11.7), Inches(4.4), RGBColor(0xED, 0xEB, 0xF7))
    _text(s, Inches(1.1), Inches(2.2), Inches(11.1), Inches(3.8),
          [("Drop your architecture diagram here", 20, VIOLET, True),
           (note, 16, SLATE, False)], space=8)
    return s


def closing_slide(prs, c):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, INK)
    _rect(s, 0, Inches(3.5), W, Inches(0.10), TEAL)
    _text(s, Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.2),
          [("Thank you", 48, WHITE, True)])
    _text(s, Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.8),
          [("Team " + c.get("team_name", "") + "  ·  " + c.get("idea", ""),
            18, RGBColor(0xC7, 0xC2, 0xF0), False)])


def build(content_path, out_path):
    with open(content_path, "r", encoding="utf-8") as handle:
        c = json.load(handle)

    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    title_slide(prs, c)
    team_slide(prs, c)
    problem_slide(prs, c)
    bullets_slide(prs, "What the Agent Does", c.get("what_it_does", []))
    architecture_slide(prs, c)
    bullets_slide(prs, "Live Demo: the aha moment", c.get("demo_steps", []))
    bullets_slide(prs, "Impact", c.get("impact", []))
    bullets_slide(prs, "What We Would Build Next", c.get("next_steps", []))
    closing_slide(prs, c)

    prs.save(out_path)
    print(f"Wrote {out_path} with {len(prs.slides._sldIdLst)} slides.")


if __name__ == "__main__":
    content = sys.argv[1] if len(sys.argv) > 1 else "deck-content.json"
    out = sys.argv[2] if len(sys.argv) > 2 else "demo.pptx"
    build(content, out)
